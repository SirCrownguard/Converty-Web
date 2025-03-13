import os
import subprocess
import uuid
import io
import zipfile
from flask import Flask, render_template, request, send_file, jsonify, after_this_request, abort
from pdf2image import convert_from_path
from pptx import Presentation

app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
CONVERTED_FOLDER = "converted_files"
for folder in [UPLOAD_FOLDER, CONVERTED_FOLDER]:
    if not os.path.exists(folder):
        os.makedirs(folder)

def save_uploaded_file(file_obj, ext):
    unique_name = f"{uuid.uuid4()}{ext}"
    file_path = os.path.join(UPLOAD_FOLDER, unique_name)
    file_obj.save(file_path)
    return file_path

def convert_pdf_to_pptx(pdf_path):
    images = convert_from_path(pdf_path)
    presentation = Presentation()
    for image in images:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        iw, ih = image.size
        sw, sh = presentation.slide_width, presentation.slide_height
        ratio = iw / ih
        if ratio > (sw / sh):
            nw = sw
            nh = sw / ratio
        else:
            nh = sh
            nw = sh * ratio
        left = (sw - nw) / 2
        top = (sh - nh) / 2
        img_bytes = io.BytesIO()
        image.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        slide.shapes.add_picture(img_bytes, left, top, width=nw, height=nh)
    unique_name = f"{uuid.uuid4()}.pptx"
    output_path = os.path.join(CONVERTED_FOLDER, unique_name)
    presentation.save(output_path)
    return unique_name

def convert_pptx_to_pdf(pptx_path):
    output_folder = os.path.abspath(CONVERTED_FOLDER)
    subprocess.run(
        ["soffice", "--headless", "--convert-to", "pdf", pptx_path, "--outdir", output_folder],
        check=True
    )
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    output_filename = f"{base}.pdf"
    unique_name = f"{uuid.uuid4()}.pdf"
    src = os.path.join(output_folder, output_filename)
    dst = os.path.join(output_folder, unique_name)
    os.rename(src, dst)
    return unique_name

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/result")
def result():
    return render_template("result.html")

# Yeni: Dosya yükleme endpoint'i
@app.route("/upload", methods=["POST"])
def upload():
    file = request.files.get("file")
    conversion_type = request.form.get("conversion_type")
    if not file:
        return jsonify({"error": "No file uploaded"}), 400
    # 100 MB sınırı
    if file.content_length is not None and file.content_length > 104857600:
        return jsonify({"error": "File size exceeds 100MB"}), 400
    if conversion_type == "pdf_to_pptx":
        if not file.filename.lower().endswith(".pdf"):
            return jsonify({"error": "Only PDF files allowed"}), 400
    elif conversion_type == "pptx_to_pdf":
        if not file.filename.lower().endswith(".pptx"):
            return jsonify({"error": "Only PPTX files allowed"}), 400
    else:
        return jsonify({"error": "Invalid conversion type"}), 400
    saved_path = save_uploaded_file(file, os.path.splitext(file.filename)[1])
    return jsonify({"file_id": os.path.basename(saved_path), "original_name": file.filename})

# Yeni: Tüm yüklenen dosyalar için dönüşüm endpoint'i
@app.route("/convert_all", methods=["POST"])
def convert_all():
    data = request.get_json()
    conversion_type = data.get("conversion_type")
    file_ids = data.get("file_ids")
    if not file_ids or not isinstance(file_ids, list):
        return jsonify({"error": "No files provided"}), 400
    converted_files = []
    try:
        if conversion_type == "pdf_to_pptx":
            for fid in file_ids:
                pdf_path = os.path.join(UPLOAD_FOLDER, fid)
                if not os.path.exists(pdf_path):
                    continue
                conv_filename = convert_pdf_to_pptx(pdf_path)
                converted_files.append(conv_filename)
                os.remove(pdf_path)
        elif conversion_type == "pptx_to_pdf":
            for fid in file_ids:
                pptx_path = os.path.join(UPLOAD_FOLDER, fid)
                if not os.path.exists(pptx_path):
                    continue
                conv_filename = convert_pptx_to_pdf(pptx_path)
                converted_files.append(conv_filename)
                os.remove(pptx_path)
        else:
            return jsonify({"error": "Invalid conversion type"}), 400

        if not converted_files:
            return jsonify({"error": "No files converted"}), 400

        # Eğer tek dosya ise direkt dosya; eğer birden fazlaysa ZIP oluştur.
        if len(converted_files) == 1:
            download_url = f"/download/{converted_files[0]}"
            return jsonify({"download_url": download_url})
        else:
            zip_filename = f"{uuid.uuid4()}.zip"
            zip_path = os.path.join(CONVERTED_FOLDER, zip_filename)
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for conv in converted_files:
                    conv_path = os.path.join(CONVERTED_FOLDER, conv)
                    zipf.write(conv_path, arcname=conv)
                    os.remove(conv_path)
            download_url = f"/download/{zip_filename}"
            return jsonify({"download_url": download_url})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/download/<filename>")
def download(filename):
    file_path = os.path.join(CONVERTED_FOLDER, filename)
    if not os.path.exists(file_path):
        abort(404)
    @after_this_request
    def remove_file(response):
        try:
            os.remove(file_path)
        except Exception:
            pass
        return response
    return send_file(file_path, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
