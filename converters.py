import io
from pdf2image import convert_from_path
from pptx import Presentation
import os

# Poppler path'ini sisteminize göre ayarlayın ya da environment variable olarak verin
# POPPLER_PATH = r"C:\path\to\poppler-xx.x.x\bin" # Windows örneği
POPPLER_PATH = os.environ.get("POPPLER_PATH", None)
os.environ["POPPLER_PATH"] = r"C:\Users\sagla\OneDrive\Belgeler\Release-24.08.0-0\poppler-24.08.0\Library\bin"

def pdf_to_pptx_lightweight(pdf_path, dpi=100):
    """
    PDF'i PPTX'e kaynakları verimli kullanarak dönüştürür.
    Düşük DPI ve minimum bellek kullanımı hedefler.
    """
    try:
        images = convert_from_path(pdf_path, dpi=dpi, poppler_path=POPPLER_PATH, fmt='png', thread_count=2)
    except Exception as e:
        print(f"Poppler hatası veya PDF dönüştürme hatası: {e}")
        raise ValueError(f"PDF işlenirken bir sorunla karşılaşıldı: {e}. Poppler doğru yapılandırılmış mı?")

    # Yeni bir sunum nesnesi ilklendiriliyor
    presentation = Presentation()
    slide_width_emu = presentation.slide_width
    slide_height_emu = presentation.slide_height

    for image in images:
        # Boş bir slayt düzeni ile yeni slayt ekleniyor
        slide = presentation.slides.add_slide(presentation.slide_layouts[5]) 

        img_width_px, img_height_px = image.size
        
        image_aspect_ratio = img_width_px / img_height_px
        slide_aspect_ratio = slide_width_emu / slide_height_emu

        if image_aspect_ratio > slide_aspect_ratio:
            new_width_emu = slide_width_emu
            new_height_emu = int(new_width_emu / image_aspect_ratio)
        else:
            new_height_emu = slide_height_emu
            new_width_emu = int(new_height_emu * image_aspect_ratio)

        left_emu = int((slide_width_emu - new_width_emu) / 2)
        top_emu = int((slide_height_emu - new_height_emu) / 2)
        
        img_stream = io.BytesIO()
        image.save(img_stream, format='PNG', optimize=True)
        img_stream.seek(0)

        try:
            slide.shapes.add_picture(img_stream, left_emu, top_emu, width=new_width_emu, height=new_height_emu)
        except Exception as e:
            print(f"Slayta resim eklenirken hata: {e}")
            continue
        finally:
            img_stream.close()
            image.close()

    return presentation

CONVERTERS = {
    "pdf_to_pptx": {
        "function": pdf_to_pptx_lightweight,
        "input_accept": ".pdf",
        "output_extension": ".pptx",
        "display_name_key": "converter_pdf_to_pptx" # Dil dosyası için anahtar
    },
    # Örnek:
    # "png_to_jpg": {
    #     "function": baska_bir_png_jpg_fonksiyonu,
    #     "input_accept": ".png",
    #     "output_extension": ".jpg",
    #     "display_name_key": "converter_png_to_jpg"
    # }
}